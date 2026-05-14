"""
SAP-Branded Slide Builder Library
=================================
Reusable, tested functions for creating SAP-branded PowerPoint presentations.

Usage:
    from slide_builder import SlideBuilder

    builder = SlideBuilder()
    builder.add_cover_slide("My Presentation", "Subtitle here")
    builder.add_agenda_slide(["Topic 1", "Topic 2", "Topic 3"])
    builder.add_content_slide("Key Point", ["Bullet 1", "Bullet 2"])
    builder.add_thank_you_slide()
    builder.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
from typing import List, Optional, Dict, Any, Tuple
from dataclasses import dataclass


# =============================================================================
# SAP Brand Colors
# =============================================================================

class SAPColors:
    """Official SAP brand color palette"""
    # Primary Blues
    DARK_BLUE = RGBColor(0, 42, 134)        # #002A86 - Headers, primary brand
    MEDIUM_BLUE = RGBColor(0, 112, 242)     # #0070F2 - Accents, interactive
    LIGHT_BLUE = RGBColor(27, 144, 255)     # #1B90FF - Highlights
    LIGHTER_BLUE = RGBColor(137, 209, 255)  # #89D1FF - Backgrounds
    LIGHTEST_BLUE = RGBColor(209, 239, 255) # #D1EFFF - Subtle backgrounds

    # Neutrals
    WHITE = RGBColor(255, 255, 255)         # #FFFFFF
    BLACK = RGBColor(0, 0, 0)               # #000000
    LIGHT_GRAY = RGBColor(234, 236, 238)    # #EAECEE - Dividers
    MEDIUM_GRAY = RGBColor(169, 180, 190)   # #A9B4BE - Secondary text

    # Accents (use sparingly)
    SUCCESS_GREEN = RGBColor(151, 221, 64)  # #97DD40
    WARNING_ORANGE = RGBColor(255, 201, 51) # #FFC933
    ERROR_RED = RGBColor(238, 57, 57)       # #EE3939
    HIGHLIGHT_YELLOW = RGBColor(255, 243, 184)  # #FFF3B8


# =============================================================================
# Layout Constants
# =============================================================================

class Layouts:
    """SAP Template layout indices"""
    # Covers (0-11)
    COVER_A = 0         # Title + Picture placeholder
    COVER_B = 1
    COVER_C = 2

    # Agendas (12-13)
    AGENDA_A = 12       # Title + Text placeholder
    AGENDA_B = 13

    # Dividers (14-17)
    DIVIDER_A = 14      # Center title only
    DIVIDER_B = 15
    DIVIDER_C = 16
    DIVIDER_D = 17

    # Content (18-28)
    TITLE_ONLY = 18
    TITLE_AND_TEXT = 19         # Most common - title + body
    TWO_COLUMNS = 20            # Title + 2 text columns
    THREE_COLUMNS = 21          # Title + 3 text columns
    TWO_COL_WITH_IMAGES = 22
    THREE_COL_WITH_IMAGES = 23
    FOUR_COL_WITH_IMAGES = 24
    TEXT_WITH_IMAGE = 25        # Title + text + image (1/3)
    FULL_BLEED_IMAGE = 26
    TEXT_AND_SCREENSHOT = 27
    TITLE_AND_CONTENT = 28

    # Special (29-32)
    QUOTE = 29
    QA = 30
    THANK_YOU_A = 31
    THANK_YOU_B = 32

    # Other
    BLANK = 33


# =============================================================================
# Content Constraints
# =============================================================================

class Constraints:
    """Brand guideline constraints"""
    MAX_BULLETS_PER_SLIDE = 6
    MAX_WORDS_PER_BULLET = 8
    MAX_AGENDA_ITEMS = 5
    MIN_FONT_SIZE = 10
    MAX_BULLET_DEPTH = 2
    WHITESPACE_PERCENT_MIN = 30
    WHITESPACE_PERCENT_MAX = 40


# =============================================================================
# Typography
# =============================================================================

class Typography:
    """Font specifications"""
    PRIMARY_FONT = "72 Brand"
    FALLBACK_FONTS = ["Arial", "Calibri", "Segoe UI"]

    # Sizes in points
    TITLE_SIZE = 40
    SUBTITLE_SIZE = 24
    SECTION_HEADER_SIZE = 32
    BODY_SIZE = 16
    SMALL_SIZE = 12
    IMPACT_NUMBER_SIZE = 60

    @classmethod
    def get_font(cls) -> str:
        """Return primary font (fallback handled by PowerPoint)"""
        return cls.PRIMARY_FONT


# =============================================================================
# SlideBuilder Class
# =============================================================================

class SlideBuilder:
    """
    SAP-branded PowerPoint slide builder.

    Provides high-level methods for creating presentations using
    the official SAP template layouts and brand guidelines.
    """

    def __init__(self, template_path: Optional[str] = None, clear_template_slides: bool = True):
        """
        Initialize the builder.

        Args:
            template_path: Path to SAP template file. If None, creates from scratch.
            clear_template_slides: If True, removes all existing slides from template
                                   (keeps only layouts). Default True.
        """
        self.template_path = template_path
        self._issues: List[str] = []
        self._clear_slides = clear_template_slides

        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            self._template_slide_count = len(self.prs.slides)
            # We'll handle clearing at save time to avoid zip file issues
        else:
            # Create blank presentation with SAP dimensions
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            self._template_slide_count = 0

        self._slide_count = 0

    # =========================================================================
    # Cover Slides
    # =========================================================================

    def add_cover_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        date: Optional[str] = None,
        layout_index: int = None
    ) -> 'SlideBuilder':
        """
        Add a cover/title slide.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle or presenter name
            date: Optional date string
            layout_index: Layout to use (default: custom styled cover)
        """
        # Always use our custom styled cover for consistent branding
        # Template layouts often have decorative elements that interfere
        slide = self._create_styled_cover(title, subtitle, date)

        self._slide_count += 1
        return self

    def _create_styled_cover(self, title: str, subtitle: str = None, date: str = None):
        """Create cover slide with manual styling (fallback)"""
        # Use blank layout to avoid template decorative elements
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Dark blue background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        self._set_fill(bg, SAPColors.DARK_BLUE)
        bg.line.fill.background()

        # Accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(3.8), self.prs.slide_width, Inches(0.1)
        )
        self._set_fill(accent, SAPColors.MEDIUM_BLUE)
        accent.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5))
        self._format_text(
            title_box, title,
            font_size=Pt(44), bold=True, color=SAPColors.WHITE,
            alignment=PP_ALIGN.CENTER
        )

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
            self._format_text(
                sub_box, subtitle,
                font_size=Pt(24), color=SAPColors.LIGHTER_BLUE,
                alignment=PP_ALIGN.CENTER
            )

        # Date
        if date:
            date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
            self._format_text(
                date_box, date,
                font_size=Pt(14), color=SAPColors.LIGHTER_BLUE,
                alignment=PP_ALIGN.CENTER
            )

        return slide

    # =========================================================================
    # Agenda Slides
    # =========================================================================

    def add_agenda_slide(
        self,
        items: List[str],
        title: str = "Agenda"
    ) -> 'SlideBuilder':
        """
        Add an agenda slide.

        Args:
            items: List of agenda items (max 5 enforced)
            title: Agenda slide title
        """
        # Enforce constraint
        if len(items) > Constraints.MAX_AGENDA_ITEMS:
            self._issues.append(f"Agenda truncated from {len(items)} to {Constraints.MAX_AGENDA_ITEMS} items")
            items = items[:Constraints.MAX_AGENDA_ITEMS]

        if self._has_layouts():
            slide = self._add_slide_from_layout(Layouts.AGENDA_A)
            self._fill_title(slide, title)
            self._fill_body(slide, items, numbered=True)
        else:
            slide = self._create_styled_agenda(title, items)

        self._slide_count += 1
        return self

    def _create_styled_agenda(self, title: str, items: List[str]):
        """Create agenda slide with manual styling (fallback)"""
        slide = self._create_content_base(title)

        # Numbered agenda items
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.5), Inches(5.0))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i+1}. {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = SAPColors.BLACK
            p.font.name = Typography.get_font()
            p.font.bold = True
            p.space_after = Pt(18)

        return slide

    # =========================================================================
    # Section Dividers
    # =========================================================================

    def add_section_divider(
        self,
        title: str,
        subtitle: Optional[str] = None
    ) -> 'SlideBuilder':
        """
        Add a section divider slide.

        Args:
            title: Section title
            subtitle: Optional subtitle
        """
        # Always use custom styled divider for consistent branding
        slide = self._create_styled_divider(title, subtitle)

        self._slide_count += 1
        return self

    def _create_styled_divider(self, title: str, subtitle: str = None):
        """Create divider slide with manual styling (fallback)"""
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Dark blue background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        self._set_fill(bg, SAPColors.DARK_BLUE)
        bg.line.fill.background()

        # Accent line
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.0), Inches(2), Inches(0.05)
        )
        self._set_fill(accent, SAPColors.MEDIUM_BLUE)
        accent.line.fill.background()

        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.5))
        self._format_text(
            title_box, title,
            font_size=Pt(40), bold=True, color=SAPColors.WHITE
        )

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(0.8))
            self._format_text(
                sub_box, subtitle,
                font_size=Pt(20), color=SAPColors.LIGHTER_BLUE
            )

        return slide

    # =========================================================================
    # Content Slides
    # =========================================================================

    def add_content_slide(
        self,
        title: str,
        bullets: List[str],
        auto_split: bool = True
    ) -> 'SlideBuilder':
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title
            bullets: List of bullet points
            auto_split: If True, creates multiple slides when exceeding limit
        """
        max_bullets = Constraints.MAX_BULLETS_PER_SLIDE

        if len(bullets) > max_bullets and auto_split:
            # Split into multiple slides
            chunks = [bullets[i:i+max_bullets] for i in range(0, len(bullets), max_bullets)]
            for i, chunk in enumerate(chunks):
                slide_title = title if i == 0 else f"{title} (continued)"
                self._create_single_content_slide(slide_title, chunk)
        else:
            if len(bullets) > max_bullets:
                self._issues.append(f"Slide '{title}' has {len(bullets)} bullets (max {max_bullets})")
            self._create_single_content_slide(title, bullets[:max_bullets])

        return self

    def _create_single_content_slide(self, title: str, bullets: List[str]):
        """Create a single content slide"""
        if self._has_layouts():
            slide = self._add_slide_from_layout(Layouts.TITLE_AND_TEXT)
            self._fill_title(slide, title)
            self._fill_body(slide, bullets)
        else:
            slide = self._create_styled_content(title, bullets)

        self._slide_count += 1
        return slide

    def _create_styled_content(self, title: str, bullets: List[str]):
        """Create content slide with manual styling (fallback)"""
        slide = self._create_content_base(title)

        # Bullet points
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = SAPColors.BLACK
            p.font.name = Typography.get_font()
            p.space_after = Pt(12)

        return slide

    def _create_content_base(self, title: str):
        """Create base content slide with header"""
        # Use blank layout (33) to avoid template decorative elements
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        self._set_fill(header, SAPColors.DARK_BLUE)
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        self._format_text(
            title_box, title,
            font_size=Pt(32), bold=True, color=SAPColors.WHITE
        )

        return slide

    # =========================================================================
    # Two-Column Slides
    # =========================================================================

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_bullets: List[str],
        right_title: str,
        right_bullets: List[str]
    ) -> 'SlideBuilder':
        """
        Add a two-column comparison slide.

        Args:
            title: Main slide title
            left_title: Left column header
            left_bullets: Left column bullet points
            right_title: Right column header
            right_bullets: Right column bullet points
        """
        if self._has_layouts():
            slide = self._add_slide_from_layout(Layouts.TWO_COLUMNS)
            self._fill_title(slide, title)
            # Fill placeholders for columns
            self._fill_column(slide, 10, left_title, left_bullets)
            self._fill_column(slide, 11, right_title, right_bullets)
        else:
            slide = self._create_styled_two_column(title, left_title, left_bullets, right_title, right_bullets)

        self._slide_count += 1
        return self

    def _create_styled_two_column(self, title, left_title, left_bullets, right_title, right_bullets):
        """Create two-column slide with manual styling (fallback)"""
        slide = self._create_content_base(title)

        col_width = Inches(5.5)
        left_x = Inches(0.75)
        right_x = Inches(7.0)
        top_y = Inches(1.6)

        # Left column
        self._add_column_content(slide, left_x, top_y, col_width, left_title, left_bullets)

        # Right column
        self._add_column_content(slide, right_x, top_y, col_width, right_title, right_bullets)

        return slide

    def _add_column_content(self, slide, x, y, width, col_title, bullets):
        """Add column content to slide"""
        # Column title
        title_box = slide.shapes.add_textbox(x, y, width, Inches(0.5))
        self._format_text(
            title_box, col_title,
            font_size=Pt(20), bold=True, color=SAPColors.DARK_BLUE
        )

        # Column bullets
        content_box = slide.shapes.add_textbox(x, y + Inches(0.6), width, Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = SAPColors.BLACK
            p.font.name = Typography.get_font()
            p.space_after = Pt(10)

    # =========================================================================
    # Three-Column Slides
    # =========================================================================

    def add_three_column_slide(
        self,
        title: str,
        columns: List[Tuple[str, List[str]]]
    ) -> 'SlideBuilder':
        """
        Add a three-column slide.

        Args:
            title: Main slide title
            columns: List of 3 tuples: [(col_title, bullets), ...]
        """
        if len(columns) != 3:
            self._issues.append(f"Three-column slide expects 3 columns, got {len(columns)}")
            columns = columns[:3] if len(columns) > 3 else columns + [("", [])] * (3 - len(columns))

        if self._has_layouts():
            slide = self._add_slide_from_layout(Layouts.THREE_COLUMNS)
            self._fill_title(slide, title)
        else:
            slide = self._create_styled_three_column(title, columns)

        self._slide_count += 1
        return self

    def _create_styled_three_column(self, title, columns):
        """Create three-column slide with manual styling (fallback)"""
        slide = self._create_content_base(title)

        col_width = Inches(3.8)
        positions = [Inches(0.75), Inches(4.75), Inches(8.75)]
        top_y = Inches(1.6)

        for i, (col_title, bullets) in enumerate(columns):
            self._add_column_content(slide, positions[i], top_y, col_width, col_title, bullets)

        return slide

    # =========================================================================
    # Quote Slides
    # =========================================================================

    def add_quote_slide(
        self,
        quote: str,
        attribution: str,
        title: Optional[str] = None
    ) -> 'SlideBuilder':
        """
        Add a quote/testimonial slide.

        Args:
            quote: The quote text
            attribution: Who said it (name, title, company)
            title: Optional slide title
        """
        if self._has_layouts():
            slide = self._add_slide_from_layout(Layouts.QUOTE)
            if title:
                self._fill_title(slide, title)
        else:
            slide = self._create_styled_quote(quote, attribution, title)

        self._slide_count += 1
        return self

    def _create_styled_quote(self, quote, attribution, title=None):
        """Create quote slide with manual styling (fallback)"""
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Light background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        self._set_fill(bg, SAPColors.LIGHTEST_BLUE)
        bg.line.fill.background()

        # Title if provided
        start_y = Inches(1.0)
        if title:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.8))
            self._format_text(
                title_box, title,
                font_size=Pt(32), bold=True, color=SAPColors.DARK_BLUE
            )
            start_y = Inches(1.8)

        # Quote mark
        quote_mark = slide.shapes.add_textbox(Inches(1), start_y, Inches(1), Inches(1))
        self._format_text(
            quote_mark, '"',
            font_size=Pt(72), bold=True, color=SAPColors.MEDIUM_BLUE
        )

        # Quote text
        quote_box = slide.shapes.add_textbox(Inches(1.5), start_y + Inches(0.8), Inches(10), Inches(3))
        self._format_text(
            quote_box, quote,
            font_size=Pt(24), color=SAPColors.BLACK
        )

        # Attribution
        attr_box = slide.shapes.add_textbox(Inches(1.5), start_y + Inches(4.0), Inches(10), Inches(0.8))
        self._format_text(
            attr_box, f"— {attribution}",
            font_size=Pt(18), color=SAPColors.MEDIUM_GRAY
        )

        return slide

    # =========================================================================
    # Thank You Slides
    # =========================================================================

    def add_thank_you_slide(
        self,
        title: str = "Thank You",
        contact_info: Optional[str] = None,
        next_steps: Optional[List[str]] = None
    ) -> 'SlideBuilder':
        """
        Add a closing/thank you slide.

        Args:
            title: Thank you message
            contact_info: Optional contact details
            next_steps: Optional list of next steps
        """
        # Always use custom styled thank you for consistent branding
        slide = self._create_styled_thank_you(title, contact_info, next_steps)

        self._slide_count += 1
        return self

    def _create_styled_thank_you(self, title, contact_info=None, next_steps=None):
        """Create thank you slide with manual styling (fallback)"""
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Dark blue background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        self._set_fill(bg, SAPColors.DARK_BLUE)
        bg.line.fill.background()

        # Thank you text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
        self._format_text(
            title_box, title,
            font_size=Pt(48), bold=True, color=SAPColors.WHITE,
            alignment=PP_ALIGN.CENTER
        )

        # Contact info
        if contact_info:
            contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.8))
            self._format_text(
                contact_box, contact_info,
                font_size=Pt(18), color=SAPColors.LIGHTER_BLUE,
                alignment=PP_ALIGN.CENTER
            )

        # Next steps
        if next_steps:
            steps_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7.3), Inches(1.5))
            tf = steps_box.text_frame
            tf.word_wrap = True

            for i, step in enumerate(next_steps):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"→ {step}"
                p.font.size = Pt(16)
                p.font.color.rgb = SAPColors.LIGHTER_BLUE
                p.font.name = Typography.get_font()
                p.alignment = PP_ALIGN.CENTER

        return slide

    # =========================================================================
    # Table Slides
    # =========================================================================

    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
        column_widths: Optional[List[float]] = None,
        highlight_header: bool = True
    ) -> 'SlideBuilder':
        """
        Add a slide with a data table.

        Args:
            title: Slide title
            headers: List of column headers
            rows: List of rows, each row is a list of cell values
            column_widths: Optional list of column widths in inches
            highlight_header: If True, header row has dark background
        """
        slide = self._create_content_base(title)

        # Calculate table dimensions
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row

        # Default widths if not provided
        if column_widths is None:
            total_width = 12.0
            column_widths = [total_width / num_cols] * num_cols

        table_width = sum(column_widths)
        row_height = min(0.5, 4.5 / num_rows)  # Adjust row height based on count

        # Create table
        table_left = Inches(0.65)
        table_top = Inches(1.6)
        table_height = Inches(row_height * num_rows)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            table_left, table_top,
            Inches(table_width), table_height
        )
        table = table_shape.table

        # Set column widths
        for i, width in enumerate(column_widths):
            table.columns[i].width = Inches(width)

        # Fill header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_table_cell(
                cell,
                bold=True,
                bg_color=SAPColors.DARK_BLUE if highlight_header else None,
                text_color=SAPColors.WHITE if highlight_header else SAPColors.DARK_BLUE
            )

        # Fill data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_value in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_value)
                    # Alternate row colors for readability
                    bg = SAPColors.LIGHTEST_BLUE if row_idx % 2 == 0 else SAPColors.WHITE
                    self._format_table_cell(cell, bg_color=bg)

        self._slide_count += 1
        return self

    def _format_table_cell(
        self,
        cell,
        bold: bool = False,
        bg_color: RGBColor = None,
        text_color: RGBColor = None,
        font_size: int = 12
    ):
        """Format a table cell"""
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

        tf = cell.text_frame
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = Typography.get_font()
            if bold:
                paragraph.font.bold = True
            if text_color:
                paragraph.font.color.rgb = text_color

    def add_status_table_slide(
        self,
        title: str,
        items: List[Tuple[str, str, str]],
        status_column: str = "Status"
    ) -> 'SlideBuilder':
        """
        Add a table with RAG status indicators.

        Args:
            title: Slide title
            items: List of tuples (item_name, description, status)
                   status can be: "green"/"on track", "yellow"/"at risk", "red"/"blocked"
            status_column: Header for status column
        """
        slide = self._create_content_base(title)

        # Create table
        num_rows = len(items) + 1
        table_shape = slide.shapes.add_table(
            num_rows, 3,
            Inches(0.65), Inches(1.6),
            Inches(12), Inches(min(5.0, 0.5 * num_rows))
        )
        table = table_shape.table

        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(7.5)
        table.columns[2].width = Inches(1.5)

        # Header row
        headers = ["Item", "Description", status_column]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_table_cell(cell, bold=True, bg_color=SAPColors.DARK_BLUE, text_color=SAPColors.WHITE)

        # Data rows with status colors
        status_colors = {
            "green": SAPColors.SUCCESS_GREEN,
            "on track": SAPColors.SUCCESS_GREEN,
            "yellow": SAPColors.WARNING_ORANGE,
            "at risk": SAPColors.WARNING_ORANGE,
            "red": SAPColors.ERROR_RED,
            "blocked": SAPColors.ERROR_RED,
        }

        for row_idx, (item_name, description, status) in enumerate(items):
            # Item name
            cell = table.cell(row_idx + 1, 0)
            cell.text = item_name
            self._format_table_cell(cell, bold=True)

            # Description
            cell = table.cell(row_idx + 1, 1)
            cell.text = description
            self._format_table_cell(cell)

            # Status with color
            cell = table.cell(row_idx + 1, 2)
            cell.text = status.title()
            status_color = status_colors.get(status.lower(), SAPColors.MEDIUM_GRAY)
            self._format_table_cell(cell, bold=True, bg_color=status_color, text_color=SAPColors.WHITE)

        self._slide_count += 1
        return self

    # =========================================================================
    # Chart Slides
    # =========================================================================

    def add_chart_placeholder_slide(
        self,
        title: str,
        chart_title: str,
        chart_description: str,
        insight: Optional[str] = None
    ) -> 'SlideBuilder':
        """
        Add a slide with a chart placeholder (for manual chart insertion).

        Note: python-pptx has limited chart support. This creates a placeholder
        with instructions for adding a chart manually in PowerPoint.

        Args:
            title: Slide title
            chart_title: Title for the chart area
            chart_description: Description of what data to chart
            insight: Optional key insight to highlight
        """
        slide = self._create_content_base(title)

        # Chart placeholder area
        chart_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.65), Inches(1.6),
            Inches(8.5), Inches(5.0)
        )
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = SAPColors.LIGHTEST_BLUE
        chart_box.line.color.rgb = SAPColors.MEDIUM_BLUE

        # Chart title
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.8), Inches(8.5), Inches(0.6))
        self._format_text(
            title_box, chart_title,
            font_size=Pt(18), bold=True, color=SAPColors.DARK_BLUE,
            alignment=PP_ALIGN.CENTER
        )

        # Chart description/placeholder text
        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(6.5), Inches(1.5))
        self._format_text(
            desc_box, f"[Insert chart here]\n{chart_description}",
            font_size=Pt(14), color=SAPColors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # Insight callout (if provided)
        if insight:
            insight_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(9.5), Inches(1.6),
                Inches(3.2), Inches(5.0)
            )
            insight_box.fill.solid()
            insight_box.fill.fore_color.rgb = SAPColors.HIGHLIGHT_YELLOW
            insight_box.line.fill.background()

            insight_title = slide.shapes.add_textbox(Inches(9.7), Inches(1.8), Inches(2.8), Inches(0.5))
            self._format_text(
                insight_title, "Key Insight",
                font_size=Pt(14), bold=True, color=SAPColors.DARK_BLUE
            )

            insight_text = slide.shapes.add_textbox(Inches(9.7), Inches(2.4), Inches(2.8), Inches(4.0))
            self._format_text(
                insight_text, insight,
                font_size=Pt(12), color=SAPColors.BLACK
            )

        self._slide_count += 1
        return self

    # =========================================================================
    # Image/Screenshot Slides
    # =========================================================================

    def add_image_slide(
        self,
        title: str,
        image_path: str,
        caption: Optional[str] = None,
        position: str = "center"
    ) -> 'SlideBuilder':
        """
        Add a slide with an image/screenshot.

        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional caption below the image
            position: "center", "left", or "right"
        """
        slide = self._create_content_base(title)

        if not os.path.exists(image_path):
            self._issues.append(f"Image not found: {image_path}")
            # Add placeholder
            placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            self._format_text(
                placeholder, f"[Image not found: {image_path}]",
                font_size=Pt(16), color=SAPColors.ERROR_RED,
                alignment=PP_ALIGN.CENTER
            )
        else:
            # Calculate image position
            max_width = Inches(11.0)
            max_height = Inches(4.5) if caption else Inches(5.0)

            if position == "left":
                left = Inches(0.65)
                max_width = Inches(7.0)
            elif position == "right":
                left = Inches(5.7)
                max_width = Inches(7.0)
            else:  # center
                left = Inches(1.15)

            # Add image (PowerPoint will maintain aspect ratio)
            pic = slide.shapes.add_picture(
                image_path, left, Inches(1.6),
                width=max_width
            )

            # Ensure image fits
            if pic.height > max_height:
                ratio = max_height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = max_height

            # Center horizontally if needed
            if position == "center":
                pic.left = int((self.prs.slide_width - pic.width) / 2)

        # Caption
        if caption:
            cap_box = slide.shapes.add_textbox(Inches(0.65), Inches(6.3), Inches(12), Inches(0.5))
            self._format_text(
                cap_box, caption,
                font_size=Pt(12), color=SAPColors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        self._slide_count += 1
        return self

    def add_screenshot_with_callouts_slide(
        self,
        title: str,
        image_path: str,
        callouts: List[str],
        callout_position: str = "right"
    ) -> 'SlideBuilder':
        """
        Add a slide with a screenshot and numbered callouts.

        Args:
            title: Slide title
            image_path: Path to the screenshot
            callouts: List of callout descriptions
            callout_position: "right" or "left" for callout text placement
        """
        slide = self._create_content_base(title)

        # Determine layout
        if callout_position == "right":
            img_left = Inches(0.65)
            img_width = Inches(7.5)
            callout_left = Inches(8.5)
        else:
            img_left = Inches(5.0)
            img_width = Inches(7.5)
            callout_left = Inches(0.65)

        # Add image
        if os.path.exists(image_path):
            pic = slide.shapes.add_picture(
                image_path, img_left, Inches(1.6),
                width=img_width
            )
            # Constrain height
            if pic.height > Inches(5.0):
                ratio = Inches(5.0) / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = Inches(5.0)
        else:
            self._issues.append(f"Screenshot not found: {image_path}")

        # Add callouts
        callout_box = slide.shapes.add_textbox(callout_left, Inches(1.6), Inches(4.0), Inches(5.0))
        tf = callout_box.text_frame
        tf.word_wrap = True

        for i, callout in enumerate(callouts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i+1}. {callout}"
            p.font.size = Pt(14)
            p.font.color.rgb = SAPColors.BLACK
            p.font.name = Typography.get_font()
            p.space_after = Pt(12)

        self._slide_count += 1
        return self

    # =========================================================================
    # Progress/Timeline Slides
    # =========================================================================

    def add_progress_slide(
        self,
        title: str,
        progress_items: List[Tuple[str, int, str]],
        show_percentage: bool = True
    ) -> 'SlideBuilder':
        """
        Add a slide showing progress bars.

        Args:
            title: Slide title
            progress_items: List of tuples (label, percentage, status_text)
            show_percentage: If True, show percentage number
        """
        slide = self._create_content_base(title)

        bar_height = Inches(0.3)
        bar_width = Inches(8.0)
        start_y = Inches(1.8)
        row_spacing = Inches(0.9)

        for i, (label, percentage, status) in enumerate(progress_items):
            y = start_y + Inches(i * 0.9)

            # Label
            label_box = slide.shapes.add_textbox(Inches(0.65), y, Inches(3.0), Inches(0.4))
            self._format_text(
                label_box, label,
                font_size=Pt(14), bold=True, color=SAPColors.DARK_BLUE
            )

            # Background bar (gray)
            bg_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(3.8), y + Inches(0.05),
                bar_width, bar_height
            )
            self._set_fill(bg_bar, SAPColors.LIGHT_GRAY)
            bg_bar.line.fill.background()

            # Progress bar (colored based on percentage)
            if percentage > 0:
                progress_width = Inches(8.0 * (percentage / 100))
                if percentage >= 80:
                    color = SAPColors.SUCCESS_GREEN
                elif percentage >= 50:
                    color = SAPColors.MEDIUM_BLUE
                else:
                    color = SAPColors.WARNING_ORANGE

                progress_bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(3.8), y + Inches(0.05),
                    progress_width, bar_height
                )
                self._set_fill(progress_bar, color)
                progress_bar.line.fill.background()

            # Percentage text
            if show_percentage:
                pct_box = slide.shapes.add_textbox(Inches(12.0), y, Inches(0.8), Inches(0.4))
                self._format_text(
                    pct_box, f"{percentage}%",
                    font_size=Pt(14), bold=True, color=SAPColors.DARK_BLUE,
                    alignment=PP_ALIGN.RIGHT
                )

            # Status text
            if status:
                status_box = slide.shapes.add_textbox(Inches(3.8), y + Inches(0.4), bar_width, Inches(0.3))
                self._format_text(
                    status_box, status,
                    font_size=Pt(11), color=SAPColors.MEDIUM_GRAY
                )

        self._slide_count += 1
        return self

    def add_timeline_slide(
        self,
        title: str,
        milestones: List[Tuple[str, str, str]],
        current_index: int = -1
    ) -> 'SlideBuilder':
        """
        Add a horizontal timeline slide.

        Args:
            title: Slide title
            milestones: List of tuples (date, milestone_name, status)
                       status: "completed", "current", "upcoming"
            current_index: Index of current milestone (for highlighting)
        """
        slide = self._create_content_base(title)

        num_milestones = len(milestones)
        if num_milestones == 0:
            return self

        # Timeline line
        line_y = Inches(3.5)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.65), line_y,
            Inches(12), Inches(0.05)
        )
        self._set_fill(line, SAPColors.MEDIUM_GRAY)
        line.line.fill.background()

        # Calculate spacing
        spacing = 11.5 / max(num_milestones - 1, 1)

        for i, (date, name, status) in enumerate(milestones):
            x = Inches(0.65 + i * spacing)

            # Milestone marker
            marker_size = Inches(0.3)
            if status == "completed":
                color = SAPColors.SUCCESS_GREEN
            elif status == "current" or i == current_index:
                color = SAPColors.MEDIUM_BLUE
            else:
                color = SAPColors.LIGHT_GRAY

            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.15), line_y - Inches(0.125),
                marker_size, marker_size
            )
            self._set_fill(marker, color)
            marker.line.fill.background()

            # Date (above line)
            date_box = slide.shapes.add_textbox(x - Inches(0.5), line_y - Inches(0.7), Inches(1.2), Inches(0.5))
            self._format_text(
                date_box, date,
                font_size=Pt(10), color=SAPColors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

            # Milestone name (below line)
            name_box = slide.shapes.add_textbox(x - Inches(0.7), line_y + Inches(0.3), Inches(1.5), Inches(1.0))
            text_color = SAPColors.DARK_BLUE if status == "current" or i == current_index else SAPColors.BLACK
            self._format_text(
                name_box, name,
                font_size=Pt(11), bold=(status == "current"), color=text_color,
                alignment=PP_ALIGN.CENTER
            )

        self._slide_count += 1
        return self

    # =========================================================================
    # Feature Showcase Slides (Roadmap/Demo Pattern)
    # =========================================================================

    def add_feature_showcase_slide(
        self,
        title: str,
        description: str,
        status: str = "GA",
        image_path: Optional[str] = None,
        key_points: Optional[List[str]] = None
    ) -> 'SlideBuilder':
        """
        Add a feature showcase slide (common in roadmap and product demos).

        Creates a slide with:
        - Title at top
        - Status badge (Beta/GA) in top-right
        - Large image/screenshot area
        - Description at bottom
        - Optional key points

        Args:
            title: Feature name
            description: Brief description of the feature
            status: "Beta", "GA", "POC", "Planned" - shown as badge
            image_path: Optional path to screenshot/image
            key_points: Optional list of key benefits
        """
        # Use blank layout for full control
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title at top
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.5))
        self._format_text(
            title_box, title,
            font_size=Pt(28), bold=True, color=SAPColors.DARK_BLUE
        )

        # Status badge (top-right)
        badge_colors = {
            "ga": SAPColors.SUCCESS_GREEN,
            "beta": SAPColors.WARNING_ORANGE,
            "poc": SAPColors.MEDIUM_BLUE,
            "planned": SAPColors.LIGHT_GRAY,
        }
        badge_color = badge_colors.get(status.lower(), SAPColors.MEDIUM_BLUE)

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(11.0), Inches(0.3),
            Inches(1.8), Inches(0.5)
        )
        self._set_fill(badge, badge_color)
        badge.line.fill.background()

        badge_text = slide.shapes.add_textbox(Inches(11.0), Inches(0.35), Inches(1.8), Inches(0.4))
        self._format_text(
            badge_text, status.upper(),
            font_size=Pt(14), bold=True, color=SAPColors.WHITE,
            alignment=PP_ALIGN.CENTER
        )

        # Image area (or placeholder)
        img_left = Inches(0.65)
        img_top = Inches(1.0)
        img_width = Inches(12.0)
        img_height = Inches(4.8)

        if image_path and os.path.exists(image_path):
            pic = slide.shapes.add_picture(
                image_path, img_left, img_top,
                width=img_width
            )
            # Constrain height
            if pic.height > img_height:
                ratio = img_height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = img_height
            # Center horizontally
            pic.left = int((self.prs.slide_width - pic.width) / 2)
        else:
            # Placeholder box
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                img_left, img_top,
                img_width, img_height
            )
            self._set_fill(placeholder, SAPColors.LIGHTEST_BLUE)
            placeholder.line.color.rgb = SAPColors.LIGHT_GRAY

            placeholder_text = slide.shapes.add_textbox(img_left, Inches(3.2), img_width, Inches(1))
            self._format_text(
                placeholder_text, "[Screenshot / Demo Video]",
                font_size=Pt(18), color=SAPColors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        # Description at bottom
        desc_box = slide.shapes.add_textbox(Inches(0.65), Inches(6.0), Inches(12), Inches(0.5))
        self._format_text(
            desc_box, description,
            font_size=Pt(14), color=SAPColors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # Key points (if provided, show on right side)
        if key_points:
            points_box = slide.shapes.add_textbox(Inches(9.5), Inches(1.2), Inches(3.2), Inches(4.5))
            tf = points_box.text_frame
            tf.word_wrap = True

            for i, point in enumerate(key_points[:4]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"+ {point}"
                p.font.size = Pt(11)
                p.font.color.rgb = SAPColors.DARK_BLUE
                p.font.name = Typography.get_font()
                p.space_after = Pt(8)

        self._slide_count += 1
        return self

    def add_transformation_slide(
        self,
        title: str,
        before_title: str,
        before_points: List[str],
        after_title: str,
        after_points: List[str],
        transition_text: Optional[str] = None
    ) -> 'SlideBuilder':
        """
        Add a transformation/journey slide showing Past → Future comparison.

        Common pattern for showing evolution, migration, or change initiatives.

        Args:
            title: Main slide title
            before_title: Header for "before" state (e.g., "Past", "Legacy", "Current")
            before_points: Points describing the before state
            after_title: Header for "after" state (e.g., "Future", "New", "Target")
            after_points: Points describing the after state
            transition_text: Optional text for the transition arrow/divider
        """
        # Use blank layout for full control
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title at top
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        self._format_text(
            title_box, title,
            font_size=Pt(28), bold=True, color=SAPColors.DARK_BLUE
        )

        # Divider line in center
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.55), Inches(1.5),
            Inches(0.05), Inches(5.5)
        )
        self._set_fill(divider, SAPColors.MEDIUM_GRAY)
        divider.line.fill.background()

        # "Before" column (left side)
        before_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.0), Inches(1.5),
            Inches(1.5), Inches(0.5)
        )
        self._set_fill(before_header, SAPColors.MEDIUM_GRAY)
        before_header.line.fill.background()

        before_text = slide.shapes.add_textbox(Inches(2.0), Inches(1.55), Inches(1.5), Inches(0.4))
        self._format_text(
            before_text, before_title,
            font_size=Pt(14), bold=True, color=SAPColors.WHITE,
            alignment=PP_ALIGN.CENTER
        )

        # Before points
        before_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5))
        tf = before_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(before_points[:5]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.font.color.rgb = SAPColors.MEDIUM_GRAY
            p.font.name = Typography.get_font()
            p.space_after = Pt(10)

        # "After" column (right side)
        after_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.5), Inches(1.5),
            Inches(1.5), Inches(0.5)
        )
        self._set_fill(after_header, SAPColors.SUCCESS_GREEN)
        after_header.line.fill.background()

        after_label = slide.shapes.add_textbox(Inches(9.5), Inches(1.55), Inches(1.5), Inches(0.4))
        self._format_text(
            after_label, after_title,
            font_size=Pt(14), bold=True, color=SAPColors.WHITE,
            alignment=PP_ALIGN.CENTER
        )

        # After points
        after_box = slide.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.8), Inches(4.5))
        tf = after_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(after_points[:5]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.font.color.rgb = SAPColors.DARK_BLUE
            p.font.name = Typography.get_font()
            p.font.bold = True
            p.space_after = Pt(10)

        # Transition arrow in center
        if transition_text:
            arrow_box = slide.shapes.add_textbox(Inches(5.8), Inches(3.5), Inches(1.5), Inches(0.5))
            self._format_text(
                arrow_box, f"→ {transition_text}",
                font_size=Pt(12), color=SAPColors.MEDIUM_BLUE,
                alignment=PP_ALIGN.CENTER
            )

        self._slide_count += 1
        return self

    def add_capability_cards_slide(
        self,
        title: str,
        capabilities: List[Tuple[str, str, str]],
        badge_text: Optional[str] = None
    ) -> 'SlideBuilder':
        """
        Add a capability overview slide with icon cards.

        Common pattern for product overviews showing multiple capabilities.

        Args:
            title: Slide title
            capabilities: List of tuples (icon_text, capability_name, description)
                         icon_text can be emoji or short text like "AI", "API", etc.
            badge_text: Optional badge in top-right (e.g., "By end of 2026")
        """
        # Use blank layout
        layout_idx = Layouts.BLANK if len(self.prs.slide_layouts) > Layouts.BLANK else 0
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
        self._format_text(
            title_box, title,
            font_size=Pt(28), bold=True, color=SAPColors.DARK_BLUE
        )

        # Badge (optional)
        if badge_text:
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(10.5), Inches(0.3),
                Inches(2.5), Inches(0.5)
            )
            self._set_fill(badge, SAPColors.MEDIUM_BLUE)
            badge.line.fill.background()

            badge_label = slide.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.4))
            self._format_text(
                badge_label, badge_text,
                font_size=Pt(12), bold=True, color=SAPColors.WHITE,
                alignment=PP_ALIGN.CENTER
            )

        # Calculate grid positions (2 rows x 3 columns max)
        card_width = Inches(3.8)
        card_height = Inches(2.5)
        positions = [
            (Inches(0.5), Inches(1.4)),    # Row 1
            (Inches(4.75), Inches(1.4)),
            (Inches(9.0), Inches(1.4)),
            (Inches(0.5), Inches(4.2)),    # Row 2
            (Inches(4.75), Inches(4.2)),
            (Inches(9.0), Inches(4.2)),
        ]

        for i, (icon, name, desc) in enumerate(capabilities[:6]):
            if i >= len(positions):
                break
            x, y = positions[i]

            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = SAPColors.LIGHTEST_BLUE
            card.line.color.rgb = SAPColors.LIGHT_GRAY

            # Icon circle
            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(0.2), y + Inches(0.2),
                Inches(0.6), Inches(0.6)
            )
            self._set_fill(icon_circle, SAPColors.DARK_BLUE)
            icon_circle.line.fill.background()

            icon_text = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.25), Inches(0.6), Inches(0.5))
            self._format_text(
                icon_text, icon[:2],  # Max 2 chars for icon
                font_size=Pt(14), bold=True, color=SAPColors.WHITE,
                alignment=PP_ALIGN.CENTER
            )

            # Capability name
            name_box = slide.shapes.add_textbox(x + Inches(0.9), y + Inches(0.25), card_width - Inches(1.1), Inches(0.5))
            self._format_text(
                name_box, name,
                font_size=Pt(14), bold=True, color=SAPColors.DARK_BLUE
            )

            # Divider line
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + Inches(0.2), y + Inches(0.9),
                card_width - Inches(0.4), Inches(0.02)
            )
            self._set_fill(line, SAPColors.MEDIUM_BLUE)
            line.line.fill.background()

            # Description
            desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.0), card_width - Inches(0.4), card_height - Inches(1.2))
            self._format_text(
                desc_box, desc,
                font_size=Pt(11), color=SAPColors.BLACK
            )

        self._slide_count += 1
        return self

    def add_release_phases_slide(
        self,
        title: str,
        phases: List[Tuple[str, str, str, List[str]]]
    ) -> 'SlideBuilder':
        """
        Add a release phases timeline slide (Alpha → Beta → GA pattern).

        Common for product roadmap presentations showing release phases.

        Args:
            title: Slide title
            phases: List of tuples (phase_name, date, status_color, features)
                   status_color: "gray" (alpha), "orange" (beta), "green" (ga)
        """
        slide = self._create_content_base(title)

        # Timeline line
        line_y = Inches(2.0)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.65), line_y,
            Inches(12), Inches(0.03)
        )
        self._set_fill(line, SAPColors.MEDIUM_GRAY)
        line.line.fill.background()

        # Calculate positions based on number of phases
        num_phases = min(len(phases), 4)
        spacing = 11.5 / max(num_phases - 1, 1)

        color_map = {
            "gray": SAPColors.MEDIUM_GRAY,
            "orange": SAPColors.WARNING_ORANGE,
            "green": SAPColors.SUCCESS_GREEN,
            "blue": SAPColors.MEDIUM_BLUE,
        }

        for i, (phase_name, date, status_color, features) in enumerate(phases[:4]):
            x = Inches(0.65 + i * spacing)
            color = color_map.get(status_color.lower(), SAPColors.MEDIUM_BLUE)

            # Phase badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x - Inches(0.5), line_y - Inches(0.5),
                Inches(1.2), Inches(0.4)
            )
            self._set_fill(badge, color)
            badge.line.fill.background()

            badge_text = slide.shapes.add_textbox(x - Inches(0.5), line_y - Inches(0.45), Inches(1.2), Inches(0.35))
            self._format_text(
                badge_text, phase_name,
                font_size=Pt(12), bold=True, color=SAPColors.WHITE,
                alignment=PP_ALIGN.CENTER
            )

            # Date label
            date_box = slide.shapes.add_textbox(x - Inches(0.5), line_y + Inches(0.1), Inches(1.2), Inches(0.3))
            self._format_text(
                date_box, date,
                font_size=Pt(10), color=SAPColors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

            # Feature box below
            feature_box_width = Inches(2.8)
            feature_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x - Inches(0.8), line_y + Inches(0.6),
                feature_box_width, Inches(3.5)
            )
            feature_box.fill.solid()
            feature_box.fill.fore_color.rgb = SAPColors.LIGHTEST_BLUE
            feature_box.line.color.rgb = color

            # Features list
            features_text = slide.shapes.add_textbox(
                x - Inches(0.7), line_y + Inches(0.8),
                feature_box_width - Inches(0.2), Inches(3.2)
            )
            tf = features_text.text_frame
            tf.word_wrap = True

            for j, feature in enumerate(features[:5]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {feature}"
                p.font.size = Pt(10)
                p.font.color.rgb = SAPColors.BLACK
                p.font.name = Typography.get_font()
                p.space_after = Pt(6)

        self._slide_count += 1
        return self

    def add_roadmap_table_slide(
        self,
        title: str,
        categories: List[str],
        timeframes: List[str],
        items: List[List[List[str]]]
    ) -> 'SlideBuilder':
        """
        Add a roadmap table slide with timeline columns.

        Common for showing features/capabilities across time periods.

        Args:
            title: Slide title
            categories: Row headers (e.g., ["Core Features", "AI", "Integration"])
            timeframes: Column headers (e.g., ["H1 2026", "H2 2026", "2027+"])
            items: 2D list of feature lists - items[category_idx][timeframe_idx] = [features]
        """
        slide = self._create_content_base(title)

        num_cols = len(timeframes) + 1  # +1 for category column
        num_rows = len(categories) + 1  # +1 for header row

        # Create table
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.65), Inches(1.6),
            Inches(12), Inches(5.2)
        )
        table = table_shape.table

        # Set column widths
        table.columns[0].width = Inches(2.5)  # Category column
        remaining_width = 9.5 / len(timeframes)
        for i in range(1, num_cols):
            table.columns[i].width = Inches(remaining_width)

        # Header row
        for j, timeframe in enumerate(timeframes):
            cell = table.cell(0, j + 1)
            cell.text = timeframe
            self._format_table_cell(cell, bold=True, bg_color=SAPColors.DARK_BLUE, text_color=SAPColors.WHITE)

        # Category column header (empty)
        self._format_table_cell(table.cell(0, 0), bg_color=SAPColors.DARK_BLUE)

        # Fill data
        for i, category in enumerate(categories):
            # Category cell
            cell = table.cell(i + 1, 0)
            cell.text = category
            self._format_table_cell(cell, bold=True, bg_color=SAPColors.LIGHTEST_BLUE, text_color=SAPColors.DARK_BLUE)

            # Feature cells
            for j in range(len(timeframes)):
                cell = table.cell(i + 1, j + 1)
                if i < len(items) and j < len(items[i]):
                    cell.text = "\n".join([f"• {f}" for f in items[i][j][:4]])
                self._format_table_cell(cell, font_size=10)

        self._slide_count += 1
        return self

    # =========================================================================
    # Executive Summary & Dashboard Slides (PPR Pattern)
    # =========================================================================

    def add_executive_summary_slide(
        self,
        title: str,
        sections: List[Tuple[str, List[Tuple[str, str, str]]]]
    ) -> 'SlideBuilder':
        """
        Add an executive summary dashboard slide with multiple sections.

        Common PPR pattern showing KPIs across different categories.

        Args:
            title: Slide title
            sections: List of tuples (section_name, metrics)
                     Each metric is (value, label, trend) where trend is "up"/"down"/""
        """
        slide = self._create_content_base(title)

        num_sections = min(len(sections), 4)
        section_width = 12.0 / num_sections
        start_x = Inches(0.65)

        for i, (section_name, metrics) in enumerate(sections[:4]):
            x = start_x + Inches(i * section_width)

            # Section header
            header_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, Inches(1.5),
                Inches(section_width - 0.2), Inches(0.4)
            )
            self._set_fill(header_box, SAPColors.DARK_BLUE)
            header_box.line.fill.background()

            header_text = slide.shapes.add_textbox(x, Inches(1.53), Inches(section_width - 0.2), Inches(0.35))
            self._format_text(
                header_text, section_name,
                font_size=Pt(11), bold=True, color=SAPColors.WHITE,
                alignment=PP_ALIGN.CENTER
            )

            # Metrics in this section
            metric_height = 4.5 / max(len(metrics), 1)
            for j, (value, label, trend) in enumerate(metrics[:4]):
                y = Inches(2.1 + j * metric_height)

                # Value with trend color
                trend_color = SAPColors.SUCCESS_GREEN if trend == "up" else SAPColors.ERROR_RED if trend == "down" else SAPColors.DARK_BLUE
                value_box = slide.shapes.add_textbox(x, y, Inches(section_width - 0.2), Inches(0.5))
                self._format_text(
                    value_box, value,
                    font_size=Pt(24), bold=True, color=trend_color,
                    alignment=PP_ALIGN.CENTER
                )

                # Trend arrow
                if trend:
                    arrow = " ^" if trend == "up" else " v"
                    arrow_box = slide.shapes.add_textbox(x + Inches(section_width - 0.6), y + Inches(0.1), Inches(0.3), Inches(0.3))
                    self._format_text(
                        arrow_box, arrow,
                        font_size=Pt(12), bold=True, color=trend_color
                    )

                # Label
                label_box = slide.shapes.add_textbox(x, y + Inches(0.45), Inches(section_width - 0.2), Inches(0.3))
                self._format_text(
                    label_box, label,
                    font_size=Pt(9), color=SAPColors.MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER
                )

        self._slide_count += 1
        return self

    def add_funnel_slide(
        self,
        title: str,
        stages: List[Tuple[str, str, str, str]]
    ) -> 'SlideBuilder':
        """
        Add an adoption/conversion funnel slide.

        Common PPR pattern for showing user journey or pipeline stages.

        Args:
            title: Slide title
            stages: List of tuples (stage_name, value, change, sub_text)
                   change format: "+13%" or "-5%"
        """
        slide = self._create_content_base(title)

        num_stages = min(len(stages), 5)
        stage_width = 11.5 / num_stages
        start_x = Inches(0.9)

        for i, (stage_name, value, change, sub_text) in enumerate(stages[:5]):
            x = start_x + Inches(i * stage_width)

            # Funnel box (narrower as you go down)
            box_reduction = i * 0.1
            box_width = Inches(stage_width - 0.3 - box_reduction)
            box_x = x + Inches(box_reduction / 2)

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                box_x, Inches(2.0),
                box_width, Inches(3.5)
            )
            # Color gradient from dark to light
            colors = [SAPColors.DARK_BLUE, SAPColors.MEDIUM_BLUE, SAPColors.LIGHT_BLUE,
                     SAPColors.LIGHTER_BLUE, SAPColors.LIGHTEST_BLUE]
            self._set_fill(box, colors[min(i, len(colors)-1)])
            box.line.fill.background()

            # Stage name
            name_box = slide.shapes.add_textbox(box_x, Inches(2.1), box_width, Inches(0.4))
            text_color = SAPColors.WHITE if i < 3 else SAPColors.DARK_BLUE
            self._format_text(
                name_box, stage_name,
                font_size=Pt(11), bold=True, color=text_color,
                alignment=PP_ALIGN.CENTER
            )

            # Big value
            value_box = slide.shapes.add_textbox(box_x, Inches(2.8), box_width, Inches(0.8))
            self._format_text(
                value_box, value,
                font_size=Pt(28), bold=True, color=text_color,
                alignment=PP_ALIGN.CENTER
            )

            # Change indicator
            if change:
                is_positive = change.startswith("+")
                change_color = SAPColors.SUCCESS_GREEN if is_positive else SAPColors.ERROR_RED
                change_box = slide.shapes.add_textbox(box_x, Inches(3.6), box_width, Inches(0.3))
                self._format_text(
                    change_box, change,
                    font_size=Pt(12), bold=True, color=change_color,
                    alignment=PP_ALIGN.CENTER
                )

            # Sub text
            if sub_text:
                sub_box = slide.shapes.add_textbox(box_x, Inches(4.0), box_width, Inches(0.6))
                self._format_text(
                    sub_box, sub_text,
                    font_size=Pt(9), color=text_color,
                    alignment=PP_ALIGN.CENTER
                )

            # Arrow between stages (except last)
            if i < num_stages - 1:
                arrow_box = slide.shapes.add_textbox(x + Inches(stage_width - 0.3), Inches(3.3), Inches(0.4), Inches(0.5))
                self._format_text(
                    arrow_box, "→",
                    font_size=Pt(20), color=SAPColors.MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER
                )

        self._slide_count += 1
        return self

    def add_okr_slide(
        self,
        title: str,
        objectives: List[Tuple[str, List[Tuple[str, int, str]]]]
    ) -> 'SlideBuilder':
        """
        Add an OKR (Objectives and Key Results) tracking slide.

        Args:
            title: Slide title
            objectives: List of tuples (objective_name, key_results)
                       Each key_result is (kr_text, progress_pct, status)
                       status: "on_track", "at_risk", "missed"
        """
        slide = self._create_content_base(title)

        num_objectives = min(len(objectives), 3)
        obj_height = 5.0 / num_objectives

        for i, (obj_name, key_results) in enumerate(objectives[:3]):
            y = Inches(1.6 + i * obj_height)

            # Objective header
            obj_box = slide.shapes.add_textbox(Inches(0.65), y, Inches(12), Inches(0.4))
            self._format_text(
                obj_box, f"O{i+1}: {obj_name}",
                font_size=Pt(14), bold=True, color=SAPColors.DARK_BLUE
            )

            # Key results
            for j, (kr_text, progress, status) in enumerate(key_results[:3]):
                kr_y = y + Inches(0.5 + j * 0.5)

                # Status color
                status_colors = {
                    "on_track": SAPColors.SUCCESS_GREEN,
                    "at_risk": SAPColors.WARNING_ORANGE,
                    "missed": SAPColors.ERROR_RED,
                }
                color = status_colors.get(status, SAPColors.MEDIUM_GRAY)

                # KR text
                kr_label = slide.shapes.add_textbox(Inches(0.85), kr_y, Inches(6), Inches(0.35))
                self._format_text(
                    kr_label, f"KR{j+1}: {kr_text}",
                    font_size=Pt(11), color=SAPColors.BLACK
                )

                # Progress bar background
                bar_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(7.0), kr_y + Inches(0.05),
                    Inches(4.0), Inches(0.25)
                )
                self._set_fill(bar_bg, SAPColors.LIGHT_GRAY)
                bar_bg.line.fill.background()

                # Progress bar fill
                if progress > 0:
                    bar_fill = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(7.0), kr_y + Inches(0.05),
                        Inches(4.0 * progress / 100), Inches(0.25)
                    )
                    self._set_fill(bar_fill, color)
                    bar_fill.line.fill.background()

                # Percentage
                pct_box = slide.shapes.add_textbox(Inches(11.2), kr_y, Inches(0.8), Inches(0.35))
                self._format_text(
                    pct_box, f"{progress}%",
                    font_size=Pt(11), bold=True, color=color,
                    alignment=PP_ALIGN.RIGHT
                )

        self._slide_count += 1
        return self

    def add_win_loss_slide(
        self,
        title: str,
        wins: List[Tuple[str, str]],
        losses: List[Tuple[str, str]],
        summary_stats: Optional[Dict[str, str]] = None
    ) -> 'SlideBuilder':
        """
        Add a win/loss analysis slide.

        Common PPR pattern for showing deal outcomes.

        Args:
            title: Slide title
            wins: List of tuples (deal_name, value) for won deals
            losses: List of tuples (deal_name, reason) for lost deals
            summary_stats: Optional dict with summary numbers like {"Win Rate": "65%", "Total Value": "€2.5M"}
        """
        slide = self._create_content_base(title)

        # Summary stats at top (if provided)
        if summary_stats:
            stat_width = 11.5 / len(summary_stats)
            for i, (stat_name, stat_value) in enumerate(summary_stats.items()):
                x = Inches(0.65 + i * stat_width)
                stat_box = slide.shapes.add_textbox(x, Inches(1.5), Inches(stat_width), Inches(0.8))
                tf = stat_box.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.text = stat_value
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = SAPColors.DARK_BLUE
                p.font.name = Typography.get_font()
                p.alignment = PP_ALIGN.CENTER

                label_box = slide.shapes.add_textbox(x, Inches(2.2), Inches(stat_width), Inches(0.3))
                self._format_text(
                    label_box, stat_name,
                    font_size=Pt(10), color=SAPColors.MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER
                )

        content_y = Inches(2.8) if summary_stats else Inches(1.6)

        # Wins column (left)
        wins_header = slide.shapes.add_textbox(Inches(0.65), content_y, Inches(5.5), Inches(0.4))
        self._format_text(
            wins_header, f"Wins ({len(wins)})",
            font_size=Pt(14), bold=True, color=SAPColors.SUCCESS_GREEN
        )

        wins_box = slide.shapes.add_textbox(Inches(0.65), content_y + Inches(0.5), Inches(5.5), Inches(3.5))
        tf = wins_box.text_frame
        tf.word_wrap = True
        for i, (deal, value) in enumerate(wins[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"+ {deal} ({value})"
            p.font.size = Pt(11)
            p.font.color.rgb = SAPColors.BLACK
            p.font.name = Typography.get_font()
            p.space_after = Pt(6)

        # Losses column (right)
        losses_header = slide.shapes.add_textbox(Inches(6.5), content_y, Inches(5.5), Inches(0.4))
        self._format_text(
            losses_header, f"Losses ({len(losses)})",
            font_size=Pt(14), bold=True, color=SAPColors.ERROR_RED
        )

        losses_box = slide.shapes.add_textbox(Inches(6.5), content_y + Inches(0.5), Inches(5.5), Inches(3.5))
        tf = losses_box.text_frame
        tf.word_wrap = True
        for i, (deal, reason) in enumerate(losses[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"- {deal}: {reason}"
            p.font.size = Pt(11)
            p.font.color.rgb = SAPColors.BLACK
            p.font.name = Typography.get_font()
            p.space_after = Pt(6)

        self._slide_count += 1
        return self

    # =========================================================================
    # Highlight/Lowlight/Action Slides (PPR Pattern)
    # =========================================================================

    def add_highlights_lowlights_slide(
        self,
        title: str,
        highlights: List[str],
        lowlights: List[str],
        actions: Optional[List[str]] = None
    ) -> 'SlideBuilder':
        """
        Add a highlights and lowlights slide (common in PPRs and status updates).

        Creates a two or three-column layout with:
        - Highlights (green accent) on the left
        - Lowlights (red accent) in the middle or right
        - Actions (blue accent) on the right (optional)

        Args:
            title: Slide title
            highlights: List of positive points/wins
            lowlights: List of challenges/issues
            actions: Optional list of action items
        """
        slide = self._create_content_base(title)

        # Determine layout based on whether actions are provided
        if actions:
            # Three columns
            col_width = Inches(3.8)
            positions = [Inches(0.65), Inches(4.65), Inches(8.65)]
            headers = [("Highlights", SAPColors.SUCCESS_GREEN),
                      ("Lowlights", SAPColors.ERROR_RED),
                      ("Action Items", SAPColors.MEDIUM_BLUE)]
            content_lists = [highlights, lowlights, actions]
        else:
            # Two columns
            col_width = Inches(5.8)
            positions = [Inches(0.65), Inches(6.85)]
            headers = [("Highlights", SAPColors.SUCCESS_GREEN),
                      ("Lowlights", SAPColors.ERROR_RED)]
            content_lists = [highlights, lowlights]

        header_y = Inches(1.5)
        content_y = Inches(2.0)

        for i, (pos, (header_text, accent_color), items) in enumerate(zip(positions, headers, content_lists)):
            # Column header with accent line
            header_box = slide.shapes.add_textbox(pos, header_y, col_width, Inches(0.4))
            self._format_text(
                header_box, header_text,
                font_size=Pt(16), bold=True, color=SAPColors.DARK_BLUE
            )

            # Accent line under header
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                pos, header_y + Inches(0.35),
                col_width - Inches(0.2), Inches(0.03)
            )
            self._set_fill(line, accent_color)
            line.line.fill.background()

            # Content bullets
            content_box = slide.shapes.add_textbox(pos, content_y, col_width, Inches(4.5))
            tf = content_box.text_frame
            tf.word_wrap = True

            for j, item in enumerate(items[:6]):  # Max 6 items per column
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = SAPColors.BLACK
                p.font.name = Typography.get_font()
                p.space_after = Pt(8)

        self._slide_count += 1
        return self

    def add_keeps_up_at_night_slide(
        self,
        title: str = "What Keeps Us Up at Night",
        concerns: List[Tuple[str, str]] = None,
        layout: str = "grid"
    ) -> 'SlideBuilder':
        """
        Add a 'What keeps us up at night' slide for risks/concerns (PPR pattern).

        Args:
            title: Slide title
            concerns: List of tuples (concern_title, concern_description)
            layout: "grid" (2x3 boxes) or "list" (vertical list)
        """
        slide = self._create_content_base(title)

        if not concerns:
            concerns = []

        if layout == "grid":
            # Grid layout (2 rows x 3 columns)
            box_width = Inches(3.8)
            box_height = Inches(2.2)
            positions = [
                (Inches(0.65), Inches(1.6)),    # Row 1, Col 1
                (Inches(4.65), Inches(1.6)),    # Row 1, Col 2
                (Inches(8.65), Inches(1.6)),    # Row 1, Col 3
                (Inches(0.65), Inches(4.0)),    # Row 2, Col 1
                (Inches(4.65), Inches(4.0)),    # Row 2, Col 2
                (Inches(8.65), Inches(4.0)),    # Row 2, Col 3
            ]

            for i, (concern_title, concern_desc) in enumerate(concerns[:6]):
                if i >= len(positions):
                    break
                x, y = positions[i]

                # Concern box with warning background
                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, box_width, box_height
                )
                self._set_fill(box, SAPColors.HIGHLIGHT_YELLOW)
                box.line.color.rgb = SAPColors.WARNING_ORANGE

                # Warning icon placeholder (exclamation mark)
                icon_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4))
                self._format_text(
                    icon_box, "!",
                    font_size=Pt(20), bold=True, color=SAPColors.WARNING_ORANGE
                )

                # Concern title
                title_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(0.1), box_width - Inches(0.6), Inches(0.5))
                self._format_text(
                    title_box, concern_title,
                    font_size=Pt(14), bold=True, color=SAPColors.DARK_BLUE
                )

                # Concern description
                desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), box_width - Inches(0.4), box_height - Inches(0.8))
                self._format_text(
                    desc_box, concern_desc,
                    font_size=Pt(11), color=SAPColors.BLACK
                )

        else:  # list layout
            content_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.6), Inches(12), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, (concern_title, concern_desc) in enumerate(concerns[:5]):
                # Title paragraph
                p_title = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p_title.text = f"! {concern_title}"
                p_title.font.size = Pt(16)
                p_title.font.bold = True
                p_title.font.color.rgb = SAPColors.WARNING_ORANGE
                p_title.font.name = Typography.get_font()
                p_title.space_after = Pt(4)

                # Description paragraph
                p_desc = tf.add_paragraph()
                p_desc.text = f"   {concern_desc}"
                p_desc.font.size = Pt(12)
                p_desc.font.color.rgb = SAPColors.BLACK
                p_desc.font.name = Typography.get_font()
                p_desc.space_after = Pt(16)

        self._slide_count += 1
        return self

    def add_ask_slide(
        self,
        title: str = "Our Ask",
        asks: List[Tuple[str, str, str]] = None
    ) -> 'SlideBuilder':
        """
        Add an 'Ask' slide for requests/needs (common in executive presentations).

        Args:
            title: Slide title
            asks: List of tuples (ask_title, description, owner/deadline)
        """
        slide = self._create_content_base(title)

        if not asks:
            asks = []

        # Create table-like layout for asks
        start_y = Inches(1.6)
        row_height = Inches(1.0)

        for i, (ask_title, description, owner) in enumerate(asks[:5]):
            y = start_y + Inches(i * 1.1)

            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.65), y + Inches(0.2),
                Inches(0.5), Inches(0.5)
            )
            self._set_fill(circle, SAPColors.MEDIUM_BLUE)
            circle.line.fill.background()

            num_box = slide.shapes.add_textbox(Inches(0.65), y + Inches(0.25), Inches(0.5), Inches(0.4))
            self._format_text(
                num_box, str(i + 1),
                font_size=Pt(16), bold=True, color=SAPColors.WHITE,
                alignment=PP_ALIGN.CENTER
            )

            # Ask title
            title_box = slide.shapes.add_textbox(Inches(1.4), y + Inches(0.1), Inches(8.0), Inches(0.4))
            self._format_text(
                title_box, ask_title,
                font_size=Pt(16), bold=True, color=SAPColors.DARK_BLUE
            )

            # Description
            desc_box = slide.shapes.add_textbox(Inches(1.4), y + Inches(0.5), Inches(8.0), Inches(0.5))
            self._format_text(
                desc_box, description,
                font_size=Pt(12), color=SAPColors.BLACK
            )

            # Owner/deadline (right aligned)
            owner_box = slide.shapes.add_textbox(Inches(10.0), y + Inches(0.3), Inches(2.5), Inches(0.4))
            self._format_text(
                owner_box, owner,
                font_size=Pt(11), color=SAPColors.MEDIUM_GRAY,
                alignment=PP_ALIGN.RIGHT
            )

        self._slide_count += 1
        return self

    # =========================================================================
    # Impact/Metrics Slides
    # =========================================================================

    def add_metrics_slide(
        self,
        title: str,
        metrics: List[Tuple[str, str, str]]
    ) -> 'SlideBuilder':
        """
        Add a metrics/KPI slide with big numbers.

        Args:
            title: Slide title
            metrics: List of tuples: [(value, label, trend), ...]
                     trend can be "up", "down", or ""
        """
        slide = self._create_content_base(title)

        num_metrics = len(metrics)
        # Calculate proper spacing - leave margins and gaps between metrics
        total_width = 12.0  # Usable width in inches
        gap = 0.3  # Gap between metric boxes
        col_width = (total_width - (gap * (num_metrics - 1))) / num_metrics
        start_x = Inches(0.65)

        for i, (value, label, trend) in enumerate(metrics):
            x = start_x + Inches(i * (col_width + gap))
            self._add_metric_box(slide, x, Inches(2.2), Inches(col_width), value, label, trend)

        self._slide_count += 1
        return self

    def _add_metric_box(self, slide, x, y, width, value, label, trend):
        """Add a metric display box"""
        # Value - big number
        value_box = slide.shapes.add_textbox(x, y, width, Inches(1.2))
        color = SAPColors.SUCCESS_GREEN if trend == "up" else SAPColors.ERROR_RED if trend == "down" else SAPColors.DARK_BLUE
        self._format_text(
            value_box, value,
            font_size=Pt(54), bold=True, color=color,
            alignment=PP_ALIGN.CENTER
        )

        # Trend indicator - right after value
        if trend:
            arrow = " ^" if trend == "up" else " v" if trend == "down" else ""
            if arrow:
                trend_box = slide.shapes.add_textbox(x, y + Inches(1.1), width, Inches(0.4))
                trend_color = SAPColors.SUCCESS_GREEN if trend == "up" else SAPColors.ERROR_RED
                self._format_text(
                    trend_box, arrow,
                    font_size=Pt(20), bold=True, color=trend_color,
                    alignment=PP_ALIGN.CENTER
                )

        # Label - below value and trend
        label_y = y + Inches(1.5) if trend else y + Inches(1.2)
        label_box = slide.shapes.add_textbox(x, label_y, width, Inches(0.8))
        self._format_text(
            label_box, label,
            font_size=Pt(14), color=SAPColors.MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

    # =========================================================================
    # Helper Methods
    # =========================================================================

    def _has_layouts(self) -> bool:
        """Check if template has usable layouts"""
        return len(self.prs.slide_layouts) > 30

    def _add_slide_from_layout(self, layout_index: int):
        """Add a slide using template layout"""
        if layout_index < len(self.prs.slide_layouts):
            layout = self.prs.slide_layouts[layout_index]
            return self.prs.slides.add_slide(layout)
        else:
            # Fallback to blank
            return self.prs.slides.add_slide(self.prs.slide_layouts[0])

    def _fill_title(self, slide, text: str):
        """Fill the title placeholder"""
        if slide.shapes.title:
            slide.shapes.title.text = text
        else:
            # Find title placeholder by type
            for shape in slide.shapes:
                if shape.is_placeholder:
                    if shape.placeholder_format.type == 1:  # TITLE
                        shape.text = text
                        break

    def _fill_center_title(self, slide, text: str):
        """Fill center title placeholder (for dividers)"""
        for shape in slide.shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type == 3:  # CENTER_TITLE
                    shape.text = text
                    return
        # Fallback to regular title
        self._fill_title(slide, text)

    def _fill_body(self, slide, items: List[str], numbered: bool = False):
        """Fill body placeholder with bullet points"""
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 2:  # BODY
                tf = shape.text_frame
                for i, item in enumerate(items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"{i+1}. {item}" if numbered else item
                    p.level = 0
                break

    def _fill_column(self, slide, placeholder_idx: int, col_title: str, bullets: List[str]):
        """Fill a column placeholder"""
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == placeholder_idx:
                tf = shape.text_frame
                # First paragraph is column title
                tf.paragraphs[0].text = col_title
                tf.paragraphs[0].font.bold = True
                # Add bullets
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                break

    def _set_fill(self, shape, color: RGBColor):
        """Set solid fill color for a shape"""
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    def _format_text(
        self,
        textbox,
        text: str,
        font_size: Pt = None,
        bold: bool = False,
        color: RGBColor = None,
        alignment: PP_ALIGN = None
    ):
        """Format text in a textbox"""
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text

        if font_size:
            p.font.size = font_size
        if bold:
            p.font.bold = True
        if color:
            p.font.color.rgb = color
        if alignment:
            p.alignment = alignment

        p.font.name = Typography.get_font()

    # =========================================================================
    # Output Methods
    # =========================================================================

    def save(self, output_path: str) -> str:
        """
        Save the presentation.

        Args:
            output_path: Path to save the .pptx file

        Returns:
            Path to saved file
        """
        import tempfile
        import shutil

        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)

        # If we need to clear template slides, save to temp, reload, and delete old slides
        if self._clear_slides and self._template_slide_count > 0:
            # Save to temp file first
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                temp_path = tmp.name

            self.prs.save(temp_path)

            # Reload and delete template slides (they're at the beginning)
            from pptx import Presentation as ReloadPresentation
            prs_clean = ReloadPresentation(temp_path)

            # Delete the first N slides (template slides)
            # We need to delete from the XML
            xml_slides = prs_clean.slides._sldIdLst
            # Get the sldId elements to remove (first N)
            slides_to_remove = list(xml_slides)[:self._template_slide_count]
            for sldId in slides_to_remove:
                xml_slides.remove(sldId)

            # Save clean version
            prs_clean.save(output_path)

            # Clean up temp file
            os.unlink(temp_path)
        else:
            self.prs.save(output_path)

        return output_path

    def get_issues(self) -> List[str]:
        """Get list of any issues encountered during building"""
        return self._issues

    def get_slide_count(self) -> int:
        """Get number of slides created"""
        return self._slide_count


# =============================================================================
# Convenience Function
# =============================================================================

def create_presentation(template_path: Optional[str] = None) -> SlideBuilder:
    """
    Create a new SlideBuilder instance.

    Args:
        template_path: Optional path to SAP template

    Returns:
        SlideBuilder instance ready for use
    """
    return SlideBuilder(template_path)
