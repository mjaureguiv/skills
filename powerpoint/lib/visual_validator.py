"""
Visual Slide Validator
======================
Exports slides as images and provides visual inspection capabilities.

Requires: Pillow, python-pptx
Optional: pdf2image (for PDF export)
"""

import os
import tempfile
import subprocess
from typing import List, Tuple, Optional
from dataclasses import dataclass


@dataclass
class SlideVisualCheck:
    """Result of visual inspection for a slide"""
    slide_number: int
    image_path: str
    issues: List[str]
    passed: bool


def export_slides_as_images(
    pptx_path: str,
    output_dir: str,
    format: str = "png",
    dpi: int = 150
) -> List[str]:
    """
    Export PowerPoint slides as individual images.

    Uses COM automation on Windows (requires PowerPoint),
    LibreOffice if available, otherwise falls back to text preview.

    Args:
        pptx_path: Path to the PowerPoint file
        output_dir: Directory to save images
        format: Image format (png, jpg)
        dpi: Resolution for export

    Returns:
        List of paths to exported images
    """
    os.makedirs(output_dir, exist_ok=True)

    # Try Windows COM automation first (most accurate)
    try:
        return _export_with_com(pptx_path, output_dir, format)
    except Exception as e:
        print(f"COM export failed: {e}")

    # Try LibreOffice (cross-platform)
    try:
        return _export_with_libreoffice(pptx_path, output_dir, format)
    except Exception as e:
        print(f"LibreOffice export failed: {e}")

    # Fallback: use python-pptx to create simple thumbnails
    print("Using text-only fallback preview")
    return _export_with_pptx_fallback(pptx_path, output_dir)


def _export_with_libreoffice(pptx_path: str, output_dir: str, format: str) -> List[str]:
    """Export using LibreOffice command line"""
    # First convert to PDF
    pdf_path = os.path.join(output_dir, "temp_presentation.pdf")

    # Try soffice command
    result = subprocess.run([
        "soffice", "--headless", "--convert-to", "pdf",
        "--outdir", output_dir, pptx_path
    ], capture_output=True, timeout=120)

    if result.returncode != 0:
        raise Exception("LibreOffice conversion failed")

    # Then convert PDF pages to images using pdf2image or similar
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(pdf_path)

        image_paths = []
        for i, img in enumerate(images):
            img_path = os.path.join(output_dir, f"slide_{i+1:02d}.{format}")
            img.save(img_path, format.upper())
            image_paths.append(img_path)

        os.unlink(pdf_path)
        return image_paths
    except ImportError:
        raise Exception("pdf2image not available")


def _export_with_com(pptx_path: str, output_dir: str, format: str) -> List[str]:
    """Export using Windows COM automation (requires PowerPoint installed)"""
    import win32com.client
    import pythoncom

    # Initialize COM
    pythoncom.CoInitialize()

    powerpoint = None
    presentation = None

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # PowerPoint must be visible for Export to work
        powerpoint.Visible = True

        presentation = powerpoint.Presentations.Open(
            os.path.abspath(pptx_path),
            ReadOnly=True,
            WithWindow=False
        )

        image_paths = []
        for i, slide in enumerate(presentation.Slides):
            img_path = os.path.join(os.path.abspath(output_dir), f"slide_{i+1:02d}.{format}")
            slide.Export(img_path, format.upper())
            image_paths.append(img_path)

        return image_paths
    finally:
        if presentation:
            presentation.Close()
        if powerpoint:
            powerpoint.Quit()
        pythoncom.CoUninitialize()


def _export_with_pptx_fallback(pptx_path: str, output_dir: str) -> List[str]:
    """
    Fallback: Create simple slide previews using python-pptx.
    This doesn't render the actual slide but extracts text/structure info.
    """
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(pptx_path)
    image_paths = []

    # Create simple text-based preview images
    for i, slide in enumerate(prs.slides):
        # Create a simple preview image
        img = Image.new('RGB', (1280, 720), color='white')
        draw = ImageDraw.Draw(img)

        # Try to get a font
        try:
            font = ImageFont.truetype("arial.ttf", 20)
            font_small = ImageFont.truetype("arial.ttf", 14)
        except:
            font = ImageFont.load_default()
            font_small = font

        # Draw slide number
        draw.text((20, 20), f"Slide {i+1}", fill='black', font=font)

        # Extract and draw text content
        y_pos = 60
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()[:100]
                if text:
                    draw.text((20, y_pos), text, fill='gray', font=font_small)
                    y_pos += 25
                    if y_pos > 680:
                        break

        # Add note that this is a text-only preview
        draw.text((20, 690), "[Text-only preview - install LibreOffice for full rendering]",
                  fill='red', font=font_small)

        img_path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
        img.save(img_path)
        image_paths.append(img_path)

    return image_paths


def create_visual_review_report(
    pptx_path: str,
    output_dir: str = None
) -> Tuple[str, List[str]]:
    """
    Create a visual review of a presentation.

    Exports all slides as images and creates a summary report.

    Args:
        pptx_path: Path to the PowerPoint file
        output_dir: Directory for output (uses temp if not specified)

    Returns:
        Tuple of (report_text, list_of_image_paths)
    """
    if output_dir is None:
        output_dir = tempfile.mkdtemp(prefix="ppt_visual_review_")

    # Export slides
    image_paths = export_slides_as_images(pptx_path, output_dir)

    # Create report
    report_lines = [
        "# Visual Slide Review",
        "",
        f"Presentation: {os.path.basename(pptx_path)}",
        f"Total slides: {len(image_paths)}",
        f"Images exported to: {output_dir}",
        "",
        "## Slides for Review",
        ""
    ]

    for i, img_path in enumerate(image_paths, 1):
        report_lines.append(f"### Slide {i}")
        report_lines.append(f"![Slide {i}]({img_path})")
        report_lines.append("")
        report_lines.append("**Visual Check:**")
        report_lines.append("- [ ] No overlapping elements")
        report_lines.append("- [ ] Text is readable")
        report_lines.append("- [ ] Proper alignment")
        report_lines.append("- [ ] Consistent spacing")
        report_lines.append("- [ ] Colors look correct")
        report_lines.append("")

    report = "\n".join(report_lines)

    return report, image_paths
