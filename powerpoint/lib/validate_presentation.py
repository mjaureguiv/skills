"""
PowerPoint Presentation Validator
=================================
Validates generated presentations against SAP brand guidelines.

Usage:
    from validate_presentation import PresentationValidator

    validator = PresentationValidator("output.pptx")
    issues = validator.validate_all()

    if issues:
        print("Issues found:")
        for issue in issues:
            print(f"  - {issue}")
    else:
        print("All checks passed!")
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from typing import List, Dict, Set, Tuple, Optional
from dataclasses import dataclass
from enum import Enum


class Severity(Enum):
    """Issue severity levels"""
    CRITICAL = "critical"
    HIGH = "high"
    MEDIUM = "medium"
    LOW = "low"


@dataclass
class ValidationIssue:
    """Represents a validation issue"""
    slide_index: int
    check_id: str
    message: str
    severity: Severity
    element: Optional[str] = None

    def __str__(self):
        location = f"Slide {self.slide_index}"
        if self.element:
            location += f" ({self.element})"
        return f"[{self.severity.value.upper()}] {location}: {self.message}"


class PresentationValidator:
    """
    Validates PowerPoint presentations against SAP brand guidelines.
    """

    # Approved SAP colors (RGB tuples)
    APPROVED_COLORS = {
        (0, 42, 134),       # Dark Blue
        (0, 112, 242),      # Medium Blue
        (27, 144, 255),     # Light Blue
        (137, 209, 255),    # Lighter Blue
        (209, 239, 255),    # Lightest Blue
        (255, 255, 255),    # White
        (0, 0, 0),          # Black
        (234, 236, 238),    # Light Gray
        (169, 180, 190),    # Medium Gray
        (151, 221, 64),     # Success Green
        (255, 201, 51),     # Warning Orange
        (238, 57, 57),      # Error Red
        (255, 243, 184),    # Highlight Yellow
        (223, 18, 120),     # Magenta (charts)
    }

    # Approved fonts (including theme font placeholders)
    APPROVED_FONTS = {
        "72 Brand", "72", "72 Light", "72 Bold",
        "Arial", "Arial Bold", "Arial Narrow",
        "Calibri", "Calibri Light", "Calibri Bold",
        "Segoe UI", "Segoe UI Light", "Segoe UI Bold",
        # Theme font placeholders (used by templates)
        "+mn-lt", "+mn-ea", "+mn-cs",  # Minor (body) fonts
        "+mj-lt", "+mj-ea", "+mj-cs",  # Major (heading) fonts
    }

    # Constraints
    MAX_BULLETS_PER_SLIDE = 6
    MAX_WORDS_PER_BULLET = 8
    MIN_FONT_SIZE = 10
    MAX_FILE_SIZE_MB = 20

    def __init__(self, pptx_path: str, skip_template_slides: int = 0):
        """
        Initialize validator with a presentation file.

        Args:
            pptx_path: Path to the PowerPoint file to validate
            skip_template_slides: Number of template slides to skip (e.g., 20 for SAP template)
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"File not found: {pptx_path}")

        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
        self.issues: List[ValidationIssue] = []
        self.skip_template_slides = skip_template_slides

    def validate_all(self) -> List[ValidationIssue]:
        """
        Run all validation checks.

        Returns:
            List of ValidationIssue objects
        """
        self.issues = []

        # Run all checks
        self._check_colors()
        self._check_fonts()
        self._check_font_sizes()
        self._check_bullet_counts()
        self._check_text_length()
        self._check_italics()
        self._check_file_size()
        self._check_empty_placeholders()

        # Sort by severity then slide number
        severity_order = {Severity.CRITICAL: 0, Severity.HIGH: 1, Severity.MEDIUM: 2, Severity.LOW: 3}
        self.issues.sort(key=lambda x: (severity_order[x.severity], x.slide_index))

        return self.issues

    def _check_colors(self):
        """Check all colors against approved palette"""
        slides_to_check = list(self.prs.slides)[self.skip_template_slides:]
        for slide_idx, slide in enumerate(slides_to_check, self.skip_template_slides + 1):
            colors_found = self._extract_colors_from_slide(slide)

            for color, location in colors_found:
                if color not in self.APPROVED_COLORS:
                    hex_color = f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"
                    self.issues.append(ValidationIssue(
                        slide_index=slide_idx,
                        check_id="color_compliance",
                        message=f"Unapproved color {hex_color} used",
                        severity=Severity.CRITICAL,
                        element=location
                    ))

    def _extract_colors_from_slide(self, slide) -> List[Tuple[Tuple[int, int, int], str]]:
        """Extract all colors from a slide with their locations"""
        colors = []

        for shape in slide.shapes:
            # Fill colors
            try:
                if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                    rgb = shape.fill.fore_color.rgb
                    if rgb:
                        colors.append(((rgb[0], rgb[1], rgb[2]), f"shape fill"))
            except:
                pass

            # Text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type == 1:  # RGB
                                rgb = run.font.color.rgb
                                if rgb:
                                    colors.append(((rgb[0], rgb[1], rgb[2]), f"text"))
                        except:
                            pass

        return colors

    def _check_fonts(self):
        """Check all fonts against approved list"""
        slides_to_check = list(self.prs.slides)[self.skip_template_slides:]
        for slide_idx, slide in enumerate(slides_to_check, self.skip_template_slides + 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                font_name = run.font.name
                                if font_name and font_name not in self.APPROVED_FONTS:
                                    # Check if it's a variant of an approved font
                                    base_approved = any(
                                        approved in font_name or font_name in approved
                                        for approved in self.APPROVED_FONTS
                                    )
                                    if not base_approved:
                                        self.issues.append(ValidationIssue(
                                            slide_index=slide_idx,
                                            check_id="font_consistency",
                                            message=f"Unapproved font '{font_name}'",
                                            severity=Severity.CRITICAL
                                        ))
                            except:
                                pass

    def _check_font_sizes(self):
        """Check minimum font size requirement"""
        slides_to_check = list(self.prs.slides)[self.skip_template_slides:]
        for slide_idx, slide in enumerate(slides_to_check, self.skip_template_slides + 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.size:
                                    size_pt = run.font.size.pt
                                    if size_pt < self.MIN_FONT_SIZE:
                                        text_preview = run.text[:20] + "..." if len(run.text) > 20 else run.text
                                        self.issues.append(ValidationIssue(
                                            slide_index=slide_idx,
                                            check_id="min_font_size",
                                            message=f"Font too small ({size_pt:.0f}pt < {self.MIN_FONT_SIZE}pt)",
                                            severity=Severity.HIGH,
                                            element=f"'{text_preview}'"
                                        ))
                            except:
                                pass

    def _check_bullet_counts(self):
        """Check bullet point counts per slide"""
        slides_to_check = list(self.prs.slides)[self.skip_template_slides:]
        for slide_idx, slide in enumerate(slides_to_check, self.skip_template_slides + 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    # Count non-empty paragraphs that look like bullets
                    bullet_count = sum(
                        1 for p in tf.paragraphs
                        if p.text.strip() and (p.text.strip().startswith('•') or p.level is not None)
                    )
                    if bullet_count > self.MAX_BULLETS_PER_SLIDE:
                        self.issues.append(ValidationIssue(
                            slide_index=slide_idx,
                            check_id="max_bullets",
                            message=f"Too many bullets ({bullet_count} > {self.MAX_BULLETS_PER_SLIDE})",
                            severity=Severity.MEDIUM
                        ))

    def _check_text_length(self):
        """Check for overly long bullet points"""
        slides_to_check = list(self.prs.slides)[self.skip_template_slides:]
        for slide_idx, slide in enumerate(slides_to_check, self.skip_template_slides + 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and not text.startswith(('•', '→', '-')):
                            continue  # Skip non-bullet text

                        # Remove bullet character for word count
                        clean_text = text.lstrip('•→- ')
                        word_count = len(clean_text.split())

                        if word_count > self.MAX_WORDS_PER_BULLET:
                            preview = clean_text[:30] + "..." if len(clean_text) > 30 else clean_text
                            self.issues.append(ValidationIssue(
                                slide_index=slide_idx,
                                check_id="bullet_length",
                                message=f"Bullet too long ({word_count} words > {self.MAX_WORDS_PER_BULLET})",
                                severity=Severity.LOW,
                                element=f"'{preview}'"
                            ))

    def _check_italics(self):
        """Check for italic text (not recommended)"""
        slides_to_check = list(self.prs.slides)[self.skip_template_slides:]
        for slide_idx, slide in enumerate(slides_to_check, self.skip_template_slides + 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.italic:
                                    text_preview = run.text[:20] + "..." if len(run.text) > 20 else run.text
                                    self.issues.append(ValidationIssue(
                                        slide_index=slide_idx,
                                        check_id="no_italics",
                                        message="Italic text found (use bold instead)",
                                        severity=Severity.MEDIUM,
                                        element=f"'{text_preview}'"
                                    ))
                            except:
                                pass

    def _check_file_size(self):
        """Check file size constraint"""
        file_size_mb = os.path.getsize(self.pptx_path) / (1024 * 1024)
        if file_size_mb > self.MAX_FILE_SIZE_MB:
            self.issues.append(ValidationIssue(
                slide_index=0,
                check_id="file_size",
                message=f"File too large ({file_size_mb:.1f}MB > {self.MAX_FILE_SIZE_MB}MB)",
                severity=Severity.LOW
            ))

    def _check_empty_placeholders(self):
        """Check for empty placeholder text"""
        placeholder_markers = [
            "click to add", "add text", "add title",
            "placeholder", "type here", "insert"
        ]

        slides_to_check = list(self.prs.slides)[self.skip_template_slides:]
        for slide_idx, slide in enumerate(slides_to_check, self.skip_template_slides + 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.lower()
                    if any(marker in text for marker in placeholder_markers):
                        self.issues.append(ValidationIssue(
                            slide_index=slide_idx,
                            check_id="empty_placeholder",
                            message="Unfilled placeholder text detected",
                            severity=Severity.HIGH
                        ))

    # =========================================================================
    # Reporting Methods
    # =========================================================================

    def get_summary(self) -> Dict:
        """
        Get a summary of validation results.

        Returns:
            Dictionary with counts by severity
        """
        if not self.issues:
            self.validate_all()

        summary = {
            "total": len(self.issues),
            "critical": sum(1 for i in self.issues if i.severity == Severity.CRITICAL),
            "high": sum(1 for i in self.issues if i.severity == Severity.HIGH),
            "medium": sum(1 for i in self.issues if i.severity == Severity.MEDIUM),
            "low": sum(1 for i in self.issues if i.severity == Severity.LOW),
            "passed": len(self.issues) == 0
        }
        return summary

    def generate_report(self) -> str:
        """
        Generate a human-readable validation report.

        Returns:
            Formatted report string
        """
        if not self.issues:
            self.validate_all()

        lines = [
            "=" * 60,
            "PRESENTATION VALIDATION REPORT",
            "=" * 60,
            f"File: {self.pptx_path}",
            f"Slides: {len(self.prs.slides)}",
            ""
        ]

        summary = self.get_summary()

        if summary["passed"]:
            lines.append("[OK] ALL CHECKS PASSED!")
        else:
            lines.append(f"Issues Found: {summary['total']}")
            lines.append(f"  - Critical: {summary['critical']}")
            lines.append(f"  - High: {summary['high']}")
            lines.append(f"  - Medium: {summary['medium']}")
            lines.append(f"  - Low: {summary['low']}")
            lines.append("")
            lines.append("-" * 60)
            lines.append("DETAILS")
            lines.append("-" * 60)

            for issue in self.issues:
                lines.append(str(issue))

        lines.append("")
        lines.append("=" * 60)

        return "\n".join(lines)

    def is_valid(self, allow_medium: bool = True, allow_low: bool = True) -> bool:
        """
        Check if presentation passes validation.

        Args:
            allow_medium: If True, medium issues don't fail validation
            allow_low: If True, low issues don't fail validation

        Returns:
            True if presentation passes, False otherwise
        """
        if not self.issues:
            self.validate_all()

        for issue in self.issues:
            if issue.severity == Severity.CRITICAL:
                return False
            if issue.severity == Severity.HIGH:
                return False
            if issue.severity == Severity.MEDIUM and not allow_medium:
                return False
            if issue.severity == Severity.LOW and not allow_low:
                return False

        return True


# =============================================================================
# Convenience Functions
# =============================================================================

def validate_presentation(pptx_path: str, strict: bool = False, skip_template_slides: int = 0) -> Tuple[bool, str]:
    """
    Validate a presentation file.

    Args:
        pptx_path: Path to the PowerPoint file
        strict: If True, medium/low issues cause failure
        skip_template_slides: Number of template slides to skip validation on

    Returns:
        Tuple of (is_valid, report_string)
    """
    validator = PresentationValidator(pptx_path, skip_template_slides=skip_template_slides)
    validator.validate_all()

    is_valid = validator.is_valid(allow_medium=not strict, allow_low=not strict)
    report = validator.generate_report()

    return is_valid, report


def quick_check(pptx_path: str) -> List[str]:
    """
    Quick validation check returning list of issue messages.

    Args:
        pptx_path: Path to the PowerPoint file

    Returns:
        List of issue messages (empty if all passed)
    """
    validator = PresentationValidator(pptx_path)
    issues = validator.validate_all()
    return [str(issue) for issue in issues]
